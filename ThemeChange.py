from enum import Enum, IntEnum
from PaintMap import PaintMap
from ShapesToMap import ShapesToMap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL

# from pptx.util import Inches
# from pptx.enum.dml import FillType, PictureFillMode


class THEME_TYPE(IntEnum):
    LavenderAroma = 0
    SereneSkies = 1
    AlluringApothecary = 2
    WellnessSpa = 3
    StarryBlue=4
    ContrastingCitrus=5
    RichReds=6

class Theme:
    all_theme = [[0xC0A9BD, 0x94A7A4, 0x64766A,0xF4F2F3],
                 [0xFBE0C3,0xFFBB98,0x7D8e95,0x344648],
                 [0xFFD5AF,0xE59A59,0x888870,0x712E1E],
                 [0xBACEC2,0xE59560,0x1D3124,0xF6F4E8],
                 [0x647295,0x9F496E,0x2B262D,0xF2EBE5],
                 [0xFAE681,0xFFA101,0xB3DEE5,0x31525B],
                 [0xF8D48A,0xD69F3A,0xC34F5A,0x541412]]
    theme=[]
    color_index=0

    def __init__(self,themeType):
        self.theme = self.all_theme[themeType]

    def colorPicker(self):
        self.color_index += 1
        self.color_index %= len(self.theme)
        return self.theme[self.color_index]

class ThemeConverter:
    theme:Theme
    input_ppt_path:str
    output_ppt_path:str

    tempShapes=[]

    def __init__(self,input_ppt_path,theme):
        self.input_ppt_path=input_ppt_path
        self.output_ppt_path="new_"+input_ppt_path
        self.theme=theme

        self.analyze_ppt()

    def analyze_ppt(self):
        presentation = Presentation(self.input_ppt_path)

        for slide in presentation.slides:
            self.processing_slide(slide)
            print("\n" + "=" * 30 + "\n")

        presentation.save(self.output_ppt_path)


    def processing_slide(self,slide):
        print(f"开始处理幻灯片：Page {slide.slide_id}")

        self.tempShapes=[] ##

        self.setBackgroudGradient(slide,self.theme.theme[0],self.theme.theme[1])

        for shape in slide.shapes:
            self.processing_shape(shape)

        if len(self.tempShapes)>6:                         #slide.slide_id!=256: #318
            self.color_variator_normal()
        else:
            map = ShapesToMap(self.tempShapes).toMap()
            tempPaintMap = PaintMap(map)
            correct_color_scheme = tempPaintMap.backCorrectColorScheme()
            self.color_variator_plus(correct_color_scheme)



    def processing_shape(self, shape, indent=1):
        print("-------------------")
        indentation = "  " * indent
        print(f"{indentation}Shape Type: {shape.shape_type}")

        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            self.processing_shape_textbox(shape,indent)
            return

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            self.processing_shape_group(shape,indent)
            return

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            self.processing_shape_picture(shape,indent)
            return

        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            self.processing_shape_chart(shape, indent)
            return

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            self.processing_shape_table(shape, indent)
            return

        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            self.processing_shape_line(shape, indent)
            return

        self.processing_shape_normal(shape, indent)
        return


    def processing_shape_group(self,shape,indent=1):
        for sub_shape in shape.shapes:
            self.processing_shape(sub_shape, indent + 1)
        return

    def processing_shape_picture(self,shape,indent=1):
        pass
        return

    def processing_shape_chart(self,shape,indent=1):

        chart=shape.chart

        for series in chart.series:
            new_color=self.theme.colorPicker()
            for point in series.points:
                self.color_variator(point.format,new_color)

        return

    def processing_shape_table(self,shape,indent=1):

        title_color=self.theme.colorPicker()
        first_row_color=self.theme.colorPicker()
        second_row_color=self.theme.colorPicker()

        table=shape.table

        row_number=-1
        for row in table.rows:
            row_number+=1

            if row_number==0:
                for cell in row.cells:
                    self.color_variator(cell, title_color)
                continue

            if row_number%2==0:
                for cell in row.cells:
                    self.color_variator(cell,first_row_color)
            else:
                for cell in row.cells:
                    self.color_variator(cell,second_row_color)
        return

    def processing_shape_line(self,shape, indent=1):
        pass
        return

    def processing_shape_textbox(self,shape, indent=1):

        print(shape.text_frame.text)
        return



    def processing_shape_normal(self,shape, indent=1):
        indentation = "  " * indent

        #不处理 带文本却没有填充 的图形
        if shape.has_text_frame:
            text = shape.text_frame.text
            print(f"{indentation}Text Content: {text}")
            if text!="" and shape.fill.type == None:
                return

        #不处理填充类型为图片的图形
        if shape.fill.type == MSO_FILL.PICTURE:#6
            return
        #不处理填充类型为背景色的图形
        if shape.fill.type== MSO_FILL.BACKGROUND:#5
            return

        if shape.fill.type == None:  # 如果图形无填充
            print(f"{indentation}Shape Color: None")  # 打印None
        else:  # 如果图形没有颜色
            self.tempShapes.append(shape)
            print("###")
            print(str(shape.shape_type)+" "+str(shape.shape_id)+" "+shape.name)
            print("###")
        return


    def color_variator_plus(self,correct_color_sheme):

        index=0
        while index<len(self.tempShapes):
            new_color=self.theme.theme[correct_color_sheme[index]-1]
            self.color_variator(self.tempShapes[index], new_color)
            index+=1

    def color_variator_normal(self):

        index=0
        while index<len(self.tempShapes):
            new_color = self.theme.colorPicker()
            self.color_variator(self.tempShapes[index], new_color)
            index+=1


    def color_variator(self,shape,color):
        shape.fill.solid()  # 设置图形填充为纯色
        r = (color >> 16) & 0xff
        g = (color >> 8) & 0xff
        b = color & 0xff
        shape.fill.fore_color.rgb = RGBColor(r, g, b)

    def setBackgroudGradient(self,slide,firstColor,SecondColor):
        fr = (firstColor >> 16) & 0xff
        fg = (firstColor >> 8) & 0xff
        fb = firstColor & 0xff

        sr = (SecondColor >> 16) & 0xff
        sg = (SecondColor >> 8) & 0xff
        sb = SecondColor & 0xff

        fill = slide.background.fill  # 获取填充对象
        fill.gradient()
        fill.gradient_stops[0].color.rgb = RGBColor(fr, fg, fb)
        fill.gradient_stops[1].color.rgb = RGBColor(sr, sg, sb)

        return



def theme_change(theme_type, file_path):
    input_ppt_path = file_path

    if theme_type == 0:
        theme = Theme(THEME_TYPE.LavenderAroma)
    elif theme_type == 1:
        theme = Theme(THEME_TYPE.SereneSkies)
    elif theme_type == 2:
        theme = Theme(THEME_TYPE.AlluringApothecary)
    elif theme_type == 3:
        theme = Theme(THEME_TYPE.WellnessSpa)
    elif theme_type == 4:
        theme = Theme(THEME_TYPE.StarryBlue)
    elif theme_type == 5:
        theme = Theme(THEME_TYPE.ContrastingCitrus)
    else:
        theme = Theme(THEME_TYPE.RichReds)

    themeConverter = ThemeConverter(input_ppt_path, theme)

    return "new_" + file_path


if __name__ == "__main__":


    input_ppt_path="XX公司计划.pptx"

    theme_change(0,input_ppt_path)
    #主题提供0~6一共7种主题







