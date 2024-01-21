# 导入python-pptx模块
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
import math

# 定义一个类，把slide中的所有shape抽象成地图区域
class ShapesToMap:
    # 定义类的属性，shapes是shape的列表，map是邻接矩阵
    shapes = []
    map = []

    # 定义类的构造函数，接受shapes作为参数，并初始化map为全零矩阵
    def __init__(self, shapes):
        self.shapes = shapes
        self.map = [[0 for _ in range(len(shapes))] for _ in range(len(shapes))]



    # 定义一个函数，判断两个shape是否重叠，返回布尔值
    def isOverlap(self, shape1, shape2):
        # 获取两个shape的左上角和右下角的坐标
        x1 = shape1.left
        y1 = shape1.top
        x2 = x1 + shape1.width
        y2 = y1 + shape1.height
        x3 = shape2.left
        y3 = shape2.top
        x4 = x3 + shape2.width
        y4 = y3 + shape2.height
        # 判断两个shape是否在x轴和y轴方向上都有重叠，如果有则返回True，否则返回False
        if x1 <= x4 and x2 >= x3 and y1 <= y4 and y2 >= y3:
            return True
        else:
            return False



    # 定义一个函数，把shapes抽象成地图的邻接矩阵
    def toMap(self):
        # 遍历shapes中的每个shape
        for i in range(len(self.shapes)):
            shape = self.shapes[i]
            for j in range(i + 1, len(self.shapes)):
                other = self.shapes[j]
                if self.isOverlap(shape, other):
                    self.map[i][j] = 1
                    self.map[j][i] = 1
        # 返回map
        return self.map

#

